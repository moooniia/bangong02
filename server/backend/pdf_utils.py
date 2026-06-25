"""轻量 PDF 处理 — 面向行政场景，优先稳定而非极致效果。"""
import os
import subprocess
import zipfile
from io import BytesIO

from docx import Document

MAX_PDF_PAGES = 50
MAX_PDF_WORD_PAGES = 80
MAX_MERGE_FILES = 20


def _open_reader(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    if reader.is_encrypted:
        try:
            reader.decrypt('')
        except Exception:
            raise ValueError('该 PDF 有密码保护，请先解密后再操作')
    return reader


def _page_count(path):
    return len(_open_reader(path).pages)


def pdf_page_count(path):
    return _page_count(path)


def check_pdf_convert_pages(path, limit=MAX_PDF_WORD_PAGES):
    n = _page_count(path)
    if n > limit:
        raise ValueError(f'PDF 共 {n} 页，请控制在 {limit} 页以内，可先拆分后再转')
    return n


def check_pdf_word_pages(path, limit=MAX_PDF_WORD_PAGES):
    return check_pdf_convert_pages(path, limit)


def check_pages(path, limit=MAX_PDF_PAGES):
    n = _page_count(path)
    if n > limit:
        raise ValueError(f'页数太多（{n} 页），请控制在 {limit} 页以内，或先拆分后再处理')
    return n


def merge_pdfs(input_paths, output_path):
    from pypdf import PdfWriter

    if len(input_paths) > MAX_MERGE_FILES:
        raise ValueError(f'一次最多合并 {MAX_MERGE_FILES} 个文件')
    if len(input_paths) < 2:
        raise ValueError('请至少选择 2 个 PDF 文件')

    writer = PdfWriter()
    for p in input_paths:
        reader = _open_reader(p)
        for page in reader.pages:
            writer.add_page(page)
    with open(output_path, 'wb') as f:
        writer.write(f)


def split_pdf(input_path, output_dir, unique_name, mode='each'):
    """mode: each=每页一个文件, half=一分为二"""
    from pypdf import PdfWriter

    reader = _open_reader(input_path)
    pages = reader.pages
    n = len(pages)
    check_pages(input_path, MAX_PDF_PAGES)

    outputs = []
    if mode == 'half':
        mid = n // 2 or 1
        ranges = [(0, mid), (mid, n)]
        for i, (start, end) in enumerate(ranges):
            writer = PdfWriter()
            for p in pages[start:end]:
                writer.add_page(p)
            out = os.path.join(output_dir, f'{unique_name}_part{i + 1}.pdf')
            with open(out, 'wb') as f:
                writer.write(f)
            outputs.append(out)
    else:
        for i, page in enumerate(pages):
            writer = PdfWriter()
            writer.add_page(page)
            out = os.path.join(output_dir, f'{unique_name}_page{i + 1}.pdf')
            with open(out, 'wb') as f:
                writer.write(f)
            outputs.append(out)

    zip_path = os.path.join(output_dir, f'{unique_name}.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for out in outputs:
            zf.write(out, os.path.basename(out))
            os.remove(out)
    return zip_path


def rotate_pdf(input_path, output_path, angle=90):
    reader = _open_reader(input_path)
    check_pages(input_path, MAX_PDF_PAGES)
    from pypdf import PdfWriter

    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)
    with open(output_path, 'wb') as f:
        writer.write(f)


def compress_pdf(input_path, output_path):
    """轻度压缩 — 无 Ghostscript 时用 pypdf 流压缩，体积降幅有限但稳定。"""
    reader = _open_reader(input_path)
    check_pages(input_path, MAX_PDF_PAGES)
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.append(reader)
    for page in writer.pages:
        page.compress_content_streams()
    with open(output_path, 'wb') as f:
        writer.write(f)


def pdf_to_images_zip(input_path, output_dir, unique_name, dpi=150):
    """PDF 转图片 — 限制 DPI 和页数保护内存。"""
    import fitz

    check_pages(input_path, MAX_PDF_PAGES)
    doc = fitz.open(input_path)
    images = []
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    for i in range(len(doc)):
        pix = doc[i].get_pixmap(matrix=mat, colorspace=fitz.csRGB)
        img_path = os.path.join(output_dir, f'{unique_name}_page{i + 1}.png')
        pix.save(img_path)
        images.append(img_path)
    doc.close()

    zip_path = os.path.join(output_dir, f'{unique_name}.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for img in images:
            zf.write(img, os.path.basename(img))
            os.remove(img)
    return zip_path


def parse_page_spec(spec, total):
    """解析页码如 1,3,5-7（1-based），返回 0-based 索引集合。"""
    if not spec or not spec.strip():
        raise ValueError('请填写要删除的页码，如 1,3,5-7')
    result = set()
    for part in spec.replace('，', ',').split(','):
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            a, b = part.split('-', 1)
            start, end = int(a), int(b)
            if start > end:
                start, end = end, start
            for p in range(start, end + 1):
                if 1 <= p <= total:
                    result.add(p - 1)
        else:
            p = int(part)
            if 1 <= p <= total:
                result.add(p - 1)
    if not result:
        raise ValueError('页码无效，请检查输入')
    return result


def delete_pdf_pages(input_path, output_path, pages_spec):
    from pypdf import PdfWriter

    reader = _open_reader(input_path)
    n = len(reader.pages)
    check_pages(input_path, MAX_PDF_PAGES)
    to_delete = parse_page_spec(pages_spec, n)
    if len(to_delete) >= n:
        raise ValueError('不能删除全部页面，至少保留一页')

    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i not in to_delete:
            writer.add_page(page)
    with open(output_path, 'wb') as f:
        writer.write(f)


def _find_cn_font():
    import glob
    patterns = [
        '/usr/share/fonts/wqy-microhei/wqy-microhei.ttc',
        '/usr/share/fonts/google-droid/DroidSansFallback.ttf',
        '/usr/share/fonts/**/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/**/wqy-microhei.ttc',
        '/usr/share/fonts/**/DejaVuSans.ttf',
    ]
    for pat in patterns:
        hits = glob.glob(pat, recursive=True) if '**' in pat else ([pat] if os.path.isfile(pat) else [])
        if hits:
            return hits[0]
    return None


def _make_watermark_tile(text, font_path):
    import io
    from PIL import Image, ImageDraw, ImageFont

    tile = Image.new('RGBA', (420, 220), (0, 0, 0, 0))
    draw = ImageDraw.Draw(tile)
    try:
        font = ImageFont.truetype(font_path, 38, index=0) if font_path else ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()
    draw.text((36, 88), text, font=font, fill=(150, 150, 150, 110))
    rotated = tile.rotate(45, expand=True)
    buf = io.BytesIO()
    rotated.save(buf, format='PNG')
    return buf.getvalue(), rotated.size


def add_pdf_watermark(input_path, output_path, text):
    import fitz

    if not text.strip():
        raise ValueError('请输入水印文字')
    check_pages(input_path, MAX_PDF_PAGES)
    img_bytes, (tw, th) = _make_watermark_tile(text, _find_cn_font())
    doc = fitz.open(input_path)
    for page in doc:
        rect = page.rect
        w, h = rect.width, rect.height
        y = -th
        while y < h + th:
            x = -tw
            while x < w + tw:
                page.insert_image(
                    fitz.Rect(x, y, x + tw, y + th),
                    stream=img_bytes,
                    overlay=True,
                )
                x += int(tw * 0.82)
            y += int(th * 0.82)
    doc.save(output_path, garbage=4, deflate=True)
    doc.close()


def encrypt_pdf(input_path, output_path, password):
    from pypdf import PdfWriter

    if not password:
        raise ValueError('请设置密码')
    reader = _open_reader(input_path)
    check_pages(input_path, MAX_PDF_PAGES)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_path, 'wb') as f:
        writer.write(f)


def decrypt_pdf(input_path, output_path, password):
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_path)
    if not reader.is_encrypted:
        raise ValueError('该 PDF 没有密码，无需解密')
    if not reader.decrypt(password):
        raise ValueError('密码错误，请重试')
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, 'wb') as f:
        writer.write(f)


def pdf_to_grayscale(input_path, output_path):
    import fitz

    check_pages(input_path, MAX_PDF_PAGES)
    doc = fitz.open(input_path)
    for page in doc:
        pix = page.get_pixmap(colorspace=fitz.csGRAY)
        page.clean_contents()
        page.insert_image(page.rect, pixmap=pix)
    doc.save(output_path)
    doc.close()


def extract_pdf_images_zip(input_path, output_dir, unique_name):
    import fitz

    check_pages(input_path, MAX_PDF_PAGES)
    doc = fitz.open(input_path)
    images = []
    for i in range(len(doc)):
        for j, img in enumerate(doc.get_page_images(i)):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.n - pix.alpha > 3:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            out = os.path.join(output_dir, f'{unique_name}_p{i + 1}_{j + 1}.png')
            pix.save(out)
            images.append(out)
    doc.close()
    if not images:
        raise ValueError('这个 PDF 里没有找到可提取的图片')
    zip_path = os.path.join(output_dir, f'{unique_name}.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for img in images:
            zf.write(img, os.path.basename(img))
            os.remove(img)
    return zip_path


def pdf_layout_to_docx(pdf_path, output_path):
    """优先 pdf2docx（电子版保留表格，扫描件按页嵌图），失败再降级纯文字。"""
    try:
        from pdf2docx import Converter

        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        if os.path.getsize(output_path) > 1024:
            return 'layout'
    except Exception:
        pass

    pdf_text_to_docx(pdf_path, output_path)
    return 'text'


def _clean_table_cell(value):
    if value is None:
        return ''
    return str(value).replace('\n', ' ').strip()


_TABLE_HEADER_SETS = [
    ['时间', '任务名称', '你需要做什么', '要准备的资料'],
    ['类别', '事项', '截止时间', '需准备的资料'],
    ['频率', '事项', '需准备的资料'],
]


def _fix_cn_spaces(text):
    import re
    if not text:
        return ''
    text = re.sub(r'(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])', '', text)
    text = re.sub(r'([，、；：])\s+', r'\1', text)
    return re.sub(r'\s{2,}', ' ', text).strip()


def _join_line_words(words):
    ordered = sorted(words, key=lambda w: w['x0'])
    out = ''
    for w in ordered:
        t = w['text']
        if not out:
            out = t
        elif out[-1].isascii() and t.isascii():
            out += ' ' + t
        else:
            out += t
    return _fix_cn_spaces(out.strip())


def _cluster_page_words(page, y_tol=4):
    from collections import defaultdict

    lines = defaultdict(list)
    for w in page.extract_words(x_tolerance=2):
        lines[round(w['top'] / y_tol) * y_tol].append(w)
    return lines


def _detect_table_header(words):
    texts = {w['text'] for w in words}
    for labels in _TABLE_HEADER_SETS:
        if all(label in texts for label in labels):
            anchors = []
            for label in labels:
                for w in sorted(words, key=lambda x: x['x0']):
                    if w['text'] == label:
                        anchors.append(w['x0'])
                        break
            if len(anchors) == len(labels):
                bounds = [0]
                bounds.extend((anchors[i] + anchors[i + 1]) / 2 for i in range(len(anchors) - 1))
                bounds.append(9999)
                return labels, bounds
    return None, None


def _assign_row_columns(words, bounds, ncols):
    buckets = [[] for _ in range(ncols)]
    for w in sorted(words, key=lambda x: x['x0']):
        idx = ncols - 1
        for i in range(len(bounds) - 1):
            if bounds[i] <= w['x0'] < bounds[i + 1]:
                idx = min(i, ncols - 1)
                break
        buckets[idx].append(w)
    return [_join_line_words(bucket) for bucket in buckets]


def _is_noise_row(row):
    text = ''.join(row).strip()
    if not text:
        return True
    if text.startswith(('一、', '二、', '三、', '四、', '五、', '□', '注：')):
        return True
    if text.startswith(('安全类资料', '平安类资料', '项目级资料')):
        return True
    if 'F:\\' in text or 'F:/' in text:
        return True
    if row[0] and len(row[0]) > 35 and not any(row[1:]):
        return True
    return False


def _section_title_from_line(line):
    for prefix in ('一、', '二、', '三、', '四、', '五、'):
        if line.startswith(prefix):
            title = line.split('□')[0].strip()
            for ch in ('/', '\\', '*', '?', ':', '[', ']'):
                title = title.replace(ch, ' ')
            return title[:31] or '数据'
    return None


def _normalize_row(row, ncols=4):
    row = [_fix_cn_spaces(_clean_table_cell(c)) for c in row]
    while len(row) < ncols:
        row.append('')
    return row[:ncols]


def _extract_pdf_table_sections(pdf_path):
    """按表头/章节拆成多个表格块，便于分 sheet 输出。"""
    import pdfplumber

    sections = []
    current = None
    pending_title = None

    def start_section(title, labels, bounds, ncols):
        nonlocal current, pending_title
        if current and current['rows']:
            sections.append(current)
        use_title = (title or pending_title or '数据')[:31]
        pending_title = None
        current = {
            'title': use_title,
            'ncols': ncols,
            'bounds': bounds,
            'header': list(labels),
            'rows': [_normalize_row(labels, max(ncols, 4))],
        }

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            line_map = _cluster_page_words(page)
            for key in sorted(line_map):
                words = line_map[key]
                line = _join_line_words(words)
                section_title = _section_title_from_line(line)
                if section_title:
                    if current and len(current['rows']) > 1:
                        sections.append(current)
                        current = None
                    pending_title = section_title
                    continue

                labels, new_bounds = _detect_table_header(words)
                if labels:
                    title = (current['title'] if current else None) or pending_title or '数据'
                    start_section(title, labels, new_bounds, len(labels))
                    continue

                if current is None or current['bounds'] is None:
                    continue
                if line.startswith('□'):
                    continue

                row = _normalize_row(
                    _assign_row_columns(words, current['bounds'], current['ncols']),
                    max(current['ncols'], 4),
                )
                if _is_noise_row(row):
                    continue
                if current['rows'] and not row[0] and any(row[1:]):
                    prev = current['rows'][-1]
                    for i in range(len(row)):
                        if row[i]:
                            prev[i] = _fix_cn_spaces(f'{prev[i]} {row[i]}'.strip())
                elif any(row):
                    current['rows'].append(row)

    if current and current['rows']:
        sections.append(current)
    return sections


def _extract_pdf_table_rows(pdf_path):
    rows = []
    for sec in _extract_pdf_table_sections(pdf_path):
        if len(sec['rows']) > 1:
            rows.extend(sec['rows'])
        elif sec['rows']:
            rows.append(sec['rows'][0])
    return rows


def pdf_has_extractable_tables(pdf_path, min_rows=3):
    sections = _extract_pdf_table_sections(pdf_path)
    data_rows = sum(max(0, len(s['rows']) - 1) for s in sections)
    return data_rows >= min_rows


def _merge_same_cells(ws, col_idx, start_row, end_row):
    from openpyxl.styles import Alignment

    if end_row <= start_row:
        return
    val = ws.cell(start_row, col_idx).value
    if not val:
        return
    ws.merge_cells(
        start_row=start_row, start_column=col_idx,
        end_row=end_row, end_column=col_idx,
    )
    cell = ws.cell(start_row, col_idx)
    cell.alignment = Alignment(wrap_text=True, vertical='center')


_CELL_TEXT_COLORS = {
    '安全': 'C62828',
    '平安': '1565C0',
}

def _write_table_sheet(ws, rows, ncols):
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    header_labels = set()
    for labels in _TABLE_HEADER_SETS:
        header_labels.update(labels)

    thin = Side(style='thin', color='B0B0B0')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill('solid', fgColor='1F4E79')
    body_font = Font(name='等线', size=11)
    header_font = Font(name='等线', size=11, bold=True, color='FFFFFF')
    widths = [10, 28, 32, 28]

    merge_ranges = []
    merge_start = None
    merge_val = None

    for r_idx, row in enumerate(rows, 1):
        is_header = bool(row and row[0] in header_labels)
        for c_idx in range(ncols):
            val = row[c_idx] if c_idx < len(row) else ''
            cell = ws.cell(r_idx, c_idx + 1, val or None)
            cell.border = border
            cell.alignment = Alignment(wrap_text=True, vertical='top')
            if is_header:
                cell.font = header_font
                cell.fill = header_fill
            else:
                text_color = _CELL_TEXT_COLORS.get(val) if c_idx == 0 else None
                cell.font = Font(name='等线', size=11, color=text_color, bold=bool(text_color)
                                 ) if text_color else body_font

        if not is_header and row and row[0]:
            if row[0] == merge_val:
                if merge_start is None:
                    merge_start = r_idx - 1
            else:
                if merge_start and r_idx - 1 > merge_start:
                    merge_ranges.append((merge_start, r_idx - 1))
                merge_val = row[0]
                merge_start = r_idx
        elif merge_start and r_idx - 1 > merge_start:
            merge_ranges.append((merge_start, r_idx - 1))
            merge_start = None
            merge_val = None

    if merge_start and len(rows) > merge_start:
        merge_ranges.append((merge_start, len(rows)))

    for start, end in merge_ranges:
        _merge_same_cells(ws, 1, start, end)

    if len(rows) > 1:
        ws.freeze_panes = 'A2'
    for i, w in enumerate(widths[:ncols], 1):
        ws.column_dimensions[chr(ord('A') + i - 1)].width = w


def _parse_md_table_sections(pages_md, _vo=None):
    """从 Volc OCR 输出中解析表格（支持 HTML <table> 和 pipe 格式）。"""
    import re
    sections = []

    for md in pages_md:
        # HTML 表格优先
        if '<table' in md.lower() and _vo is not None:
            html_blocks = re.findall(r'<table[^>]*>.*?</table>', md, re.DOTALL | re.IGNORECASE)
            if not html_blocks:
                html_blocks = [md]
            for html in html_blocks:
                try:
                    parsed = _vo._parse_html_table_rows(html)
                    if not parsed:
                        continue
                    rows = [[cell['text'] for cell in row] for row in parsed]
                    ncols = max(len(r) for r in rows) if rows else 0
                    if ncols >= 2 and len(rows) >= 2:
                        rows = [(r + [''] * ncols)[:ncols] for r in rows]
                        sections.append({
                            'title': '',
                            'ncols': ncols,
                            'rows': rows,
                            'bounds': None,
                            'header': rows[0],
                        })
                except Exception:
                    pass
            continue

        # pipe 格式表格
        current_rows = []
        header_row = None
        for line in md.split('\n'):
            stripped = line.strip()
            if stripped.startswith('|') and stripped.endswith('|') and len(stripped) > 2:
                if re.match(r'^[\|\-\s:]+$', stripped):
                    if current_rows:
                        header_row = current_rows[-1]
                    continue
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                if any(c for c in cells):
                    current_rows.append(cells)
            else:
                if len(current_rows) >= 2:
                    ncols = max(len(r) for r in current_rows)
                    rows = [(r + [''] * ncols)[:ncols] for r in current_rows]
                    sections.append({
                        'title': '',
                        'ncols': ncols,
                        'rows': rows,
                        'bounds': None,
                        'header': header_row or rows[0],
                    })
                current_rows = []
                header_row = None
        if len(current_rows) >= 2:
            ncols = max(len(r) for r in current_rows)
            rows = [(r + [''] * ncols)[:ncols] for r in current_rows]
            sections.append({
                'title': '',
                'ncols': ncols,
                'rows': rows,
                'bounds': None,
                'header': header_row or rows[0],
            })

    return sections


def _extract_pdf_doc_title(pdf_path):
    """提取第一页表格之前的标题行。"""
    import pdfplumber
    title_lines = []
    with pdfplumber.open(pdf_path) as pdf:
        line_map = _cluster_page_words(pdf.pages[0])
        for key in sorted(line_map):
            words = line_map[key]
            line = _join_line_words(words)
            if _section_title_from_line(line) or _detect_table_header(words)[0]:
                break
            if line.strip():
                title_lines.append(line)
    return title_lines


def _extract_pdf_text_sections(pdf_path):
    """提取纯文字节（无表头的节，如资料清单）。"""
    import pdfplumber
    result = []
    pending_title = None
    pending_lines = []

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            line_map = _cluster_page_words(page)
            for key in sorted(line_map):
                words = line_map[key]
                line = _join_line_words(words)
                section_title = _section_title_from_line(line)
                if section_title:
                    if pending_title and pending_lines:
                        result.append({'title': pending_title, 'lines': pending_lines})
                    pending_title = section_title
                    pending_lines = []
                    continue
                if _detect_table_header(words)[0]:
                    # 这节有表格，丢弃已积累的文字
                    pending_title = None
                    pending_lines = []
                    continue
                if pending_title and line.strip():
                    pending_lines.append(line)

    if pending_title and pending_lines:
        result.append({'title': pending_title, 'lines': pending_lines})
    return result


def _write_ocr_xlsx(ws, pages_md, _vo):
    """OCR 扫描件模式：把 HTML 表格写入 Excel，保留 colspan/rowspan 合并。
    按 text→table→text→table→... 顺序处理，文字段（标题、分节、页脚）各占一行合并。
    """
    import re
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    thin = Side(style='thin', color='B0B0B0')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    hdr_fill = PatternFill('solid', fgColor='2F75B6')
    hdr_font = Font(name='等线', size=10, bold=True, color='FFFFFF')
    bold_font = Font(name='等线', size=10, bold=True)
    body_font = Font(name='等线', size=10)
    sec_font = Font(name='等线', size=11, bold=True)

    r = 1
    max_col = 1
    text_rows = []  # 所有纯文字行（标题/分节/页脚），事后统一合并

    def _write_table(html, start_r):
        nonlocal max_col
        try:
            parsed = _vo._parse_html_table_rows(html)
        except Exception:
            return start_r
        if not parsed:
            return start_r

        occupied = set()
        placements = []
        nrows = ncols_tbl = 0
        for ri, row in enumerate(parsed):
            ci = 0
            for cell in row:
                while (ri, ci) in occupied:
                    ci += 1
                rs = max(cell.get('rowspan', 1), 1)
                cs = max(cell.get('colspan', 1), 1)
                for rr in range(ri, ri + rs):
                    for cc in range(ci, ci + cs):
                        occupied.add((rr, cc))
                placements.append((ri, ci, cell, rs, cs))
                nrows = max(nrows, ri + rs)
                ncols_tbl = max(ncols_tbl, ci + cs)
                ci += cs

        max_col = max(max_col, ncols_tbl)
        tbl_start = start_r

        for abs_ri, abs_ci, cell, rs, cs in placements:
            er = tbl_start + abs_ri
            ec = abs_ci + 1
            text = (cell.get('text') or '').strip()
            xcell = ws.cell(er, ec, text or None)
            xcell.border = border
            is_header = abs_ri == 0
            is_merged = rs > 1 or cs > 1
            is_numeric = bool(text and re.match(r'^-?\d+\.?\d*$', text))
            if is_header:
                xcell.fill = hdr_fill
                xcell.font = hdr_font
                xcell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='center')
            else:
                xcell.font = bold_font if is_merged else body_font
                xcell.alignment = Alignment(
                    wrap_text=True, vertical='center',
                    horizontal='center' if (is_merged or is_numeric) else 'left',
                )
            if is_merged:
                ws.merge_cells(
                    start_row=er, start_column=ec,
                    end_row=er + rs - 1, end_column=ec + cs - 1,
                )

        return tbl_start + nrows + 1  # +1 空行分隔

    def _write_text_block(text_block, start_r):
        cur = start_r
        clean = re.sub(r'<[^>]+>', '', text_block)  # 去残留 HTML
        for line in clean.splitlines():
            line = re.sub(r'^#+\s*', '', line).strip()
            if line:
                xcell = ws.cell(cur, 1, line)
                xcell.font = sec_font
                xcell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
                text_rows.append(cur)
                cur += 1
        return cur

    for md in pages_md:
        # 按 text/table 顺序切分页面内容
        pos = 0
        for m in re.finditer(r'<table[^>]*>.*?</table>', md, re.DOTALL | re.IGNORECASE):
            if m.start() > pos:
                r = _write_text_block(md[pos:m.start()], r)
            r = _write_table(m.group(), r)
            pos = m.end()
        # 最后一个表格之后的文字（页脚、备注等）
        if pos < len(md):
            tail = md[pos:].strip()
            if tail:
                r = _write_text_block(tail, r)

    # 列宽
    for i in range(1, max_col + 1):
        ws.column_dimensions[chr(ord('A') + i - 1)].width = 18

    # 所有纯文字行跨全列合并
    if max_col > 1:
        last_col = chr(ord('A') + max_col - 1)
        for tr in text_rows:
            try:
                ws.merge_cells(f'A{tr}:{last_col}{tr}')
            except Exception:
                pass

    if r > 2:
        ws.freeze_panes = 'A2'


def _write_single_sheet(ws, doc_title, sections, ncols, text_sections=None):
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    header_labels = set()
    for labels in _TABLE_HEADER_SETS:
        header_labels.update(labels)

    thin = Side(style='thin', color='B0B0B0')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill('solid', fgColor='1F4E79')
    body_font = Font(name='等线', size=11)
    header_font = Font(name='等线', size=11, bold=True, color='FFFFFF')
    col_letter = lambda i: chr(ord('A') + i)
    last_col = col_letter(ncols - 1)

    r = 1
    # 标题块
    title_styles = [
        Font(name='等线', size=16, bold=True, color='1F4E79'),
        Font(name='等线', size=11, italic=True, color='C62828'),
        Font(name='等线', size=10, color='666666'),
    ]
    for i, line in enumerate(doc_title):
        ws.merge_cells(f'A{r}:{last_col}{r}')
        cell = ws.cell(r, 1, line)
        cell.font = title_styles[min(i, len(title_styles) - 1)]
        cell.alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[r].height = 22 if i == 0 else 18
        r += 1
    if doc_title:
        r += 1  # 空行

    for sec in sections:
        # 节标题行
        sec_title = sec.get('title') or ''
        sec_color = 'C62828' if '安全' in sec_title else '1565C0'
        ws.merge_cells(f'A{r}:{last_col}{r}')
        cell = ws.cell(r, 1, sec_title)
        cell.font = Font(name='等线', size=12, bold=True, color=sec_color)
        cell.alignment = Alignment(vertical='center')
        ws.row_dimensions[r].height = 20
        r += 1

        # 表格行
        merge_ranges = []
        merge_start = None
        merge_val = None
        sec_start_r = r
        sec_header_set = set(sec.get('header') or [])

        for row in sec['rows']:
            is_header = bool(row and (row[0] in header_labels or
                                      (sec_header_set and row[0] in sec_header_set)))
            for c_idx in range(ncols):
                val = row[c_idx] if c_idx < len(row) else ''
                cell = ws.cell(r, c_idx + 1, val or None)
                cell.border = border
                cell.alignment = Alignment(wrap_text=True, vertical='top')
                if is_header:
                    cell.font = header_font
                    cell.fill = header_fill
                else:
                    text_color = _CELL_TEXT_COLORS.get(val) if c_idx == 0 else None
                    cell.font = Font(name='等线', size=11, color=text_color,
                                     bold=bool(text_color)) if text_color else body_font

            if not is_header and row and row[0]:
                if row[0] == merge_val:
                    if merge_start is None:
                        merge_start = r - 1
                else:
                    if merge_start and r - 1 > merge_start:
                        merge_ranges.append((merge_start, r - 1))
                    merge_val = row[0]
                    merge_start = r
            elif merge_start and r - 1 > merge_start:
                merge_ranges.append((merge_start, r - 1))
                merge_start = None
                merge_val = None
            r += 1

        if merge_start and r - 1 > merge_start:
            merge_ranges.append((merge_start, r - 1))

        for start, end in merge_ranges:
            _merge_same_cells(ws, 1, start, end)

        r += 1  # 节间空行

    # 文字节（如资料清单）
    for tsec in (text_sections or []):
        sec_title = tsec.get('title') or ''
        sec_color = 'C62828' if '安全' in sec_title else '1565C0'
        ws.merge_cells(f'A{r}:{last_col}{r}')
        cell = ws.cell(r, 1, sec_title)
        cell.font = Font(name='等线', size=12, bold=True, color=sec_color)
        cell.alignment = Alignment(vertical='center')
        ws.row_dimensions[r].height = 20
        r += 1

        for line in tsec['lines']:
            ws.merge_cells(f'A{r}:{last_col}{r}')
            cell = ws.cell(r, 1, line)
            if line.startswith('注：') or line.startswith('注:'):
                cell.font = Font(name='等线', size=10, italic=True, color='888888')
            elif '安全类' in line or '平安类' in line:
                cell.font = Font(name='等线', size=11, bold=True,
                                 color='C62828' if '安全' in line else '1565C0')
            else:
                cell.font = Font(name='等线', size=11)
            cell.alignment = Alignment(wrap_text=True, vertical='top', indent=1)
            r += 1
        r += 1  # 节间空行

    # 列宽与冻结
    widths = [10, 28, 32, 28]
    for i, w in enumerate(widths[:ncols], 1):
        ws.column_dimensions[col_letter(i - 1)].width = w
    ws.freeze_panes = 'A2'


def pdf_tables_to_xlsx(pdf_path, output_path):
    """提取 PDF 所有表格及文字节，合并到单 sheet，顶部保留文件标题。"""
    from openpyxl import Workbook

    sections = _extract_pdf_table_sections(pdf_path)
    if not sections:
        # 扫描件：走 Volc OCR 兜底，单独路径保留合并单元格
        _ocr_err = None
        try:
            import volc_ocr as _vo
            pages_md, _ = _vo._pdf_image_mode_pages(pdf_path, with_detail=False)
        except Exception as _e:
            pages_md, _vo, _ocr_err = [], None, _e
        if pages_md and _vo:
            import usage_stats as _us
            _us.bump('volc_excel')
            wb = Workbook()
            ws = wb.active
            ws.title = '表格内容'
            _write_ocr_xlsx(ws, pages_md, _vo)
            wb.save(output_path)
            return
        import logging as _log
        _log.getLogger(__name__).error('pdf_tables_to_xlsx OCR fallback failed: %s', _ocr_err)
        raise ValueError('未能提取表格内容；扫描件请先转 Word，再在 WPS 中另存为 Excel')

    doc_title = _extract_pdf_doc_title(pdf_path)
    text_sections = _extract_pdf_text_sections(pdf_path)
    ncols = max((sec['ncols'] for sec in sections), default=4)
    ncols = max(ncols, 4)

    wb = Workbook()
    ws = wb.active
    ws.title = '任务清单'
    _write_single_sheet(ws, doc_title, sections, ncols, text_sections=text_sections)

    wb.save(output_path)


def pdf_text_to_docx(pdf_path, output_path):
    """pdftotext + python-docx，版式会简化，但文字可保留。"""
    result = subprocess.run(
        ['pdftotext', pdf_path, '-'],
        capture_output=True, text=True, timeout=60,
    )
    text = result.stdout.strip()
    if len(text) < 30:
        raise ValueError('无法从 PDF 提取文字，可能是扫描件或加密文件')

    doc = Document()
    for line in text.split('\n'):
        line = line.strip()
        if line:
            doc.add_paragraph(line)
    doc.save(output_path)


def pdf_to_pptx(pdf_path, output_path, dpi=150):
    """每页渲染为 PNG 嵌入 PowerPoint，完整保留原始版式。"""
    import io as _io
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    pdf_doc = fitz.open(pdf_path)
    prs = Presentation()

    if len(pdf_doc) > 0:
        fp = pdf_doc[0]
        prs.slide_width = Emu(int(fp.rect.width * 12700))
        prs.slide_height = Emu(int(fp.rect.height * 12700))

    blank_layout = prs.slide_layouts[6]
    mat = fitz.Matrix(dpi / 72, dpi / 72)

    for page in pdf_doc:
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            _io.BytesIO(img_bytes), 0, 0,
            prs.slide_width, prs.slide_height,
        )

    pdf_doc.close()
    prs.save(output_path)


def pdf_thumbnails(input_path, max_dim=160):
    """Return list of base64 PNG data-URIs, one per page."""
    import fitz
    import base64

    doc = fitz.open(input_path)
    results = []
    for page in doc:
        rect = page.rect
        scale = max_dim / max(rect.width, rect.height)
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        data = base64.b64encode(pix.tobytes('png')).decode()
        results.append('data:image/png;base64,' + data)
    doc.close()
    return results


def pdf_page_preview(input_path, idx, max_dim=1600):
    """Render a single page at higher resolution for the editor's zoom view."""
    import fitz
    import base64

    doc = fitz.open(input_path)
    try:
        if idx < 0 or idx >= len(doc):
            raise ValueError('页码超出范围')
        page = doc[idx]
        rect = page.rect
        scale = max_dim / max(rect.width, rect.height)
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        data = base64.b64encode(pix.tobytes('png')).decode()
        return 'data:image/png;base64,' + data
    finally:
        doc.close()


def pdf_editor_export(upload_folder, pages_spec, output_path,
                      uniform_size=None, watermark_text=None,
                      grayscale=False, encrypt_password=None,
                      rotate_deg=0, compress=False):
    """
    Assemble a new PDF from pages_spec then apply optional operations.
    pages_spec: list of {"file": "server_uuid.pdf", "idx": 0_based_int}
    uniform_size: 'a4', 'a3', or None
    rotate_deg: 0, 90, 180, 270 — applied to every page
    compress: run pypdf content-stream compression on output
    """
    import fitz
    import shutil

    PAGE_SIZES = {
        'a4': fitz.paper_rect('a4'),
        'a3': fitz.paper_rect('a3'),
    }
    rotate_deg = int(rotate_deg) if rotate_deg else 0
    rotate_deg = rotate_deg % 360

    source_docs = {}
    tmp_build = output_path + '._build.pdf'
    tmp_gray  = output_path + '._gray.pdf'
    tmp_cmp   = output_path + '._cmp.pdf'

    try:
        for item in pages_spec:
            fname = item['file']
            if fname not in source_docs:
                fpath = os.path.join(upload_folder, fname)
                if not os.path.exists(fpath):
                    raise ValueError('源文件已过期，请重新上传')
                source_docs[fname] = fitz.open(fpath)

        out_doc = fitz.open()
        for item in pages_spec:
            fname = item['file']
            idx = int(item['idx'])
            src = source_docs[fname]
            if idx < 0 or idx >= len(src):
                continue
            # Per-page rotation overrides the global rotate_deg
            item_rot = int(item.get('rotate', rotate_deg)) % 360
            if uniform_size and uniform_size.lower() in PAGE_SIZES:
                base = PAGE_SIZES[uniform_size.lower()]
                bw, bh = base.width, base.height
            else:
                sr = src[idx].rect
                bw, bh = sr.width, sr.height
            if item_rot in (90, 270):
                pw, ph = bh, bw
            else:
                pw, ph = bw, bh
            new_pg = out_doc.new_page(width=pw, height=ph)
            # show_pdf_page() 的 rotate 是相对于源页面"原始未旋转内容"算的，
            # 完全不管源页面自带的 /Rotate 元数据；但缩略图/预览(get_pixmap)是会
            # 自动套用源页面自带旋转的。两者基准不一样，源页面本身带旋转时，
            # 这里如果只传 item_rot 会跟预览里看到的方向差了源页面自带的那部分角度
            # （实测：源页带180度时，预览转正后导出却整页倒了180度，正是这个原因）。
            # 把源页面自带的旋转角度叠加进去，才能让导出方向跟预览保持一致。
            effective_rotate = (item_rot + src[idx].rotation) % 360
            new_pg.show_pdf_page(new_pg.rect, src, idx, rotate=effective_rotate)

        if watermark_text and watermark_text.strip():
            img_bytes, (tw, th) = _make_watermark_tile(watermark_text.strip(), _find_cn_font())
            for page in out_doc:
                rect = page.rect
                w, h = rect.width, rect.height
                y = -th
                while y < h + th:
                    x = -tw
                    while x < w + tw:
                        page.insert_image(
                            fitz.Rect(x, y, x + tw, y + th),
                            stream=img_bytes, overlay=True,
                        )
                        x += int(tw * 0.82)
                    y += int(th * 0.82)

        for doc in source_docs.values():
            doc.close()
        source_docs = {}

        out_doc.save(tmp_build, garbage=4, deflate=True)
        out_doc.close()
        current = tmp_build

        if grayscale:
            gray_in = fitz.open(current)
            gray_out = fitz.open()
            for page in gray_in:
                pix = page.get_pixmap(colorspace=fitz.csGRAY)
                np = gray_out.new_page(width=page.rect.width, height=page.rect.height)
                np.insert_image(np.rect, pixmap=pix)
            gray_in.close()
            gray_out.save(tmp_gray)
            gray_out.close()
            if os.path.exists(current):
                os.remove(current)
            current = tmp_gray

        if compress:
            from pypdf import PdfReader, PdfWriter
            rdr = PdfReader(current)
            wtr = PdfWriter()
            wtr.append(rdr)
            for page in wtr.pages:
                page.compress_content_streams()
            with open(tmp_cmp, 'wb') as fh:
                wtr.write(fh)
            if os.path.exists(current):
                os.remove(current)
            current = tmp_cmp

        if encrypt_password:
            from pypdf import PdfReader, PdfWriter
            reader = PdfReader(current)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(encrypt_password)
            with open(output_path, 'wb') as fh:
                writer.write(fh)
            if os.path.exists(current):
                os.remove(current)
        else:
            shutil.move(current, output_path)

    finally:
        for doc in source_docs.values():
            try:
                doc.close()
            except Exception:
                pass
        for p in [tmp_build, tmp_gray, tmp_cmp]:
            if os.path.exists(p):
                os.remove(p)


def images_to_pdf_export(image_paths, rotations, output_path, uniform_size=None):
    """图片转PDF编辑器版：每张图可单独旋转，支持保持原尺寸或统一A4/A3。

    rotate 用 PIL（image_utils.rotate_image）先把图片实际转好再贴进PDF，
    不借助 PyMuPDF 自己的旋转参数——这个项目里 PyMuPDF 的旋转参数語意
    跟预期不一致已经踩过坑（pdf_editor_export 的 show_pdf_page 那次），
    图片这边干脆绕开，用已经验证过没问题的 rotate_image 来转。
    """
    if not image_paths:
        raise ValueError('请至少上传 1 张图片')
    if len(image_paths) > MAX_MERGE_FILES:
        raise ValueError(f'一次最多 {MAX_MERGE_FILES} 张图片')

    import fitz
    from image_utils import rotate_image

    PAGE_SIZES = {'a4': fitz.paper_rect('a4'), 'a3': fitz.paper_rect('a3')}
    out_doc = fitz.open()
    tmp_files = []
    try:
        for i, src_path in enumerate(image_paths):
            rotate = int(rotations[i]) % 360 if i < len(rotations) else 0
            if rotate:
                base, ext = os.path.splitext(src_path)
                tmp_path = f'{base}._rot{rotate}{ext}'
                rotate_image(src_path, tmp_path, rotate)
                tmp_files.append(tmp_path)
                img_path = tmp_path
            else:
                img_path = src_path

            pix = fitz.Pixmap(img_path)
            iw, ih = pix.width, pix.height
            pix = None

            if uniform_size and uniform_size.lower() in PAGE_SIZES:
                base_rect = PAGE_SIZES[uniform_size.lower()]
                bw, bh = base_rect.width, base_rect.height
                # 旋转90/270后图片变成横向，页面也要跟着从竖向A4/A3换成横向，
                # 不然旋转后的图片会被挤在竖向页面的宽度里，上下露白
                # （跟 pdf_editor_export 里 item_rot in (90,270) 时换宽高是同一个道理）
                if rotate in (90, 270):
                    pw, ph = bh, bw
                else:
                    pw, ph = bw, bh
                scale = min(pw / iw, ph / ih)
                dw, dh = iw * scale, ih * scale
                x0, y0 = (pw - dw) / 2, (ph - dh) / 2
                img_rect = fitz.Rect(x0, y0, x0 + dw, y0 + dh)
            else:
                pw, ph = iw * 72.0 / 96, ih * 72.0 / 96
                img_rect = fitz.Rect(0, 0, pw, ph)

            new_pg = out_doc.new_page(width=pw, height=ph)
            new_pg.insert_image(img_rect, filename=img_path)

        out_doc.save(output_path, garbage=4, deflate=True)
    finally:
        out_doc.close()
        for t in tmp_files:
            try:
                os.remove(t)
            except OSError:
                pass


def images_to_pdf(image_paths, output_path):
    if not image_paths:
        raise ValueError('请至少上传 1 张图片')
    if len(image_paths) > MAX_MERGE_FILES:
        raise ValueError(f'一次最多 {MAX_MERGE_FILES} 张图片')

    import img2pdf

    with open(output_path, 'wb') as f:
        f.write(img2pdf.convert(image_paths))